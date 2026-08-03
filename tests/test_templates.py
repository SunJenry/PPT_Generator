import os
import tempfile

from pptx import Presentation

from PPT_Generator.models import Slide
from PPT_Generator.templates.registry import TemplateRegistry


def test_registry_has_core_layouts():
    registry = TemplateRegistry()
    layouts = registry.list_layouts()
    assert "title_slide" in layouts
    assert "content" in layouts
    assert "section_divider" in layouts
    assert "closing" in layouts
    assert len(layouts) == 10


def test_title_layout_renders():
    registry = TemplateRegistry()
    layout = registry.get("title_slide")
    prs = Presentation()
    blank = prs.slide_layouts[6]
    prs_slide = prs.slides.add_slide(blank)
    slide = Slide(page_number=1, layout_id="title_slide", title="Hello", subtitle="World")
    layout.render(slide, prs_slide)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    prs.save(path)
    assert os.path.getsize(path) > 0
    os.unlink(path)
