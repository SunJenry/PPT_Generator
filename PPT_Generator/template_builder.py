"""Generate a reference template PPTX showcasing all 10 layout designs.

Usage:
    python -m PPT_Generator.template_builder [output_path]

Default output: assets/template.pptx
"""

import sys
from typing import List, Optional

from PPT_Generator.design import SLIDE_HEIGHT, SLIDE_WIDTH
from PPT_Generator.models import Slide
from PPT_Generator.templates.registry import TemplateRegistry


def _make_slide(
    page_number: int,
    layout_id: str,
    title: str,
    **kwargs,
) -> Slide:
    """Build a Slide model with all supported fields."""
    defaults: dict = {
        "subtitle": None,
        "bullets": [],
        "section_number": None,
        "section_title": None,
        "left_column": [],
        "right_column": [],
        "cards": [],
        "timeline_items": [],
        "highlight_number": None,
        "highlight_label": None,
        "quote_text": None,
        "quote_author": None,
        "closing_text": None,
        "table": None,
        "image_keyword": None,
        "image_url": None,
        "source_notes": [],
        "notes": None,
    }
    defaults.update(kwargs)
    defaults["layout_id"] = layout_id
    defaults["title"] = title
    defaults["page_number"] = page_number
    return Slide(**defaults)


def build_sample_slides() -> List[Slide]:
    """Create sample slides demonstrating every layout."""
    slides: List[Slide] = []

    # 1. Title Slide
    slides.append(_make_slide(1, "title_slide", "演示模板",
        subtitle="PPT Generator — 专业演示文稿自动生成",
        source_notes=["基于 DeepSeek 大模型 + 设计系统"],
    ))

    # 2. Section Divider
    slides.append(_make_slide(2, "section_divider", "项目背景与市场分析",
        section_number=1,
        subtitle="了解行业现状与核心驱动因素",
    ))

    # 3. Content (standard bullet)
    slides.append(_make_slide(3, "content", "核心发现概览",
        section_title="项目背景与市场分析",
        bullets=[
            "市场规模：2026年全球市场规模达到 450亿美元，年复合增长率 12.3%",
            "技术趋势：AI 原生应用成为主流，大模型驱动的自动化工具增长迅速",
            "竞争格局：头部企业占据 35% 市场份额，中小型创新企业活跃",
            "用户需求：企业客户对一站式解决方案的需求持续上升",
        ],
        source_notes=["数据来源：Gartner 2026 Market Report"],
    ))

    # 4. Two Column
    slides.append(_make_slide(4, "two_column", "优势与挑战对比",
        section_title="项目背景与市场分析",
        left_column=[
            "技术领先：自研模型推理引擎",
            "成本可控：单次生成成本 ¥0.10 以内",
            "快速迭代：3天开发周期",
            "生态兼容：支持多种 LLM 后端",
        ],
        right_column=[
            "品牌认知度不足",
            "需要持续优化模板美观度",
            "大模型输出稳定性仍需提升",
            "竞品价格战压力",
        ],
    ))

    # 5. Three Card
    slides.append(_make_slide(5, "three_card", "三大核心能力",
        cards=[
            {"title": "智能内容规划", "body": "基于主题自动生成叙事大纲，保证 25-30 页连贯内容，融合网络搜索核验事实。"},
            {"title": "专业视觉设计", "body": "10 种商务布局，统一设计系统，深蓝+金色配色方案，支持封面、时间线、对比表等。"},
            {"title": "高效低成本", "body": "平均 30 秒生成一页，单页成本 ¥0.02 以内，全自动流程无需人工干预。"},
        ],
    ))

    # 6. Section Divider
    slides.append(_make_slide(6, "section_divider", "实施路线图",
        section_number=2,
        subtitle="从规划到交付的完整路径",
    ))

    # 7. Timeline
    slides.append(_make_slide(7, "timeline", "项目关键里程碑",
        section_title="实施路线图",
        timeline_items=[
            {"date": "2026.Q3", "event": "需求分析与技术选型"},
            {"date": "2026.Q4", "event": "核心管道开发"},
            {"date": "2027.Q1", "event": "模板系统重构"},
            {"date": "2027.Q2", "event": "内部测试与优化"},
            {"date": "2027.Q3", "event": "产品发布"},
        ],
    ))

    # 8. Comparison Table
    slides.append(_make_slide(8, "comparison_table", "方案对比分析",
        section_title="实施路线图",
        table=[
            ["指标", "方案A", "方案B", "方案C"],
            ["开发周期", "3个月", "6个月", "2个月"],
            ["单页成本", "¥0.02", "¥0.15", "¥0.08"],
            ["美观度", "优秀", "良好", "一般"],
            ["可维护性", "高", "中", "低"],
            ["推荐指数", "⭐⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐"],
        ],
    ))

    # 9. Data Highlight
    slides.append(_make_slide(9, "data_highlight", "成本效益分析",
        section_title="实施路线图",
        highlight_number="¥0.02/slide",
        highlight_label="平均单页生成成本（含 LLM 调用与图片搜索）",
        notes="以 30 页演示文稿为例，总成本约 ¥0.60，相比人工制作（约 ¥200-500/页）节省 99%+",
    ))

    # 10. Quote
    slides.append(_make_slide(10, "quote", "设计理念",
        section_title="实施路线图",
        quote_text="好的演示文稿不是信息的堆砌，而是故事的讲述。每一页都应该推动叙事向前发展。",
        quote_author="PPT Generator 设计团队",
        notes="我们坚持'内容为先、视觉为辅'的原则，确保每套 PPT 都有清晰的内在叙事线。",
    ))

    # 11. Closing
    slides.append(_make_slide(11, "closing", "感谢观看",
        closing_text="期待您的反馈与建议",
        subtitle="PPT Generator v2.0 — 让每一份演示都成为精品",
    ))

    return slides


def build_template(output_path: str = "assets/template.pptx") -> str:
    """Generate the reference template .pptx file."""
    import os

    from pptx import Presentation as PptxPresentation

    slides = build_sample_slides()
    templates = TemplateRegistry()

    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    total = len(slides)
    for i, slide_model in enumerate(slides, start=1):
        prs_slide = prs.slides.add_slide(blank_layout)
        layout = templates.get(slide_model.layout_id)
        layout.render(slide_model, prs_slide, total_pages=total)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
    prs.save(output_path)
    print(f"Template saved to {output_path} ({total} slides, {len(templates.list_layouts())} layouts)")
    return output_path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "assets/template.pptx"
    build_template(out)
