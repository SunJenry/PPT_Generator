"""Centralized design system for PPT Generator.

All visual parameters and common drawing helpers live here.
Both the renderer and template builder use this module to ensure
output matches the reference template.
"""

from __future__ import annotations

from typing import List, Optional

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

# ═══════════════════════════════════════════════════════════════════
# Slide dimensions (widescreen 16:9)
# ═══════════════════════════════════════════════════════════════════

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════
# Color palette — deep navy business style
# ═══════════════════════════════════════════════════════════════════

PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)       # deep navy — titles, headers
ACCENT = RGBColor(0xC9, 0xA9, 0x62)        # warm gold — decorative highlights
SECONDARY = RGBColor(0x2E, 0x86, 0xAB)     # teal blue — secondary accents
BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF)     # white — slide background
TEXT_MAIN = RGBColor(0x2D, 0x2D, 0x2D)      # dark gray — body text
TEXT_MUTED = RGBColor(0x71, 0x71, 0x71)     # medium gray — secondary text
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF7)       # light blue-gray — card / panel bg
DIVIDER = RGBColor(0xDD, 0xDD, 0xDD)        # light gray — thin lines
HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)      # solid navy — header bar background
SUCCESS = RGBColor(0x27, 0xAE, 0x60)        # green — positive indicators
WARNING = RGBColor(0xF3, 0x9C, 0x12)        # amber — risk / warning highlights
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # white text on dark backgrounds

# ═══════════════════════════════════════════════════════════════════
# Typography
# ═══════════════════════════════════════════════════════════════════

TITLE_FONT = "Microsoft YaHei"
BODY_FONT = "Microsoft YaHei"
WESTERN_FONT = "Arial"

# Font sizes
SIZE_TITLE_COVER = Pt(44)   # cover slide title
SIZE_SUBTITLE = Pt(20)      # cover slide subtitle
SIZE_SLIDE_TITLE = Pt(34)   # regular slide title
SIZE_SECTION = Pt(30)       # section divider
SIZE_BODY = Pt(18)          # body / bullet text
SIZE_BODY_LG = Pt(20)       # large body text
SIZE_SMALL = Pt(14)         # small annotations
SIZE_CAPTION = Pt(11)       # footer / page number
SIZE_HIGHLIGHT = Pt(48)     # data highlight number
SIZE_QUOTE = Pt(28)         # quote text

# ═══════════════════════════════════════════════════════════════════
# Safe areas — guarantee content stays within slide bounds
# ═══════════════════════════════════════════════════════════════════

SAFE_LEFT = Inches(1.0)
SAFE_RIGHT = Inches(1.0)
SAFE_TOP = Inches(0.3)
SAFE_BOTTOM = Inches(0.65)

CONTENT_WIDTH = SLIDE_WIDTH - SAFE_LEFT - SAFE_RIGHT          # 11.333" (EMU)
CONTENT_WIDTH_INCHES = 11.333                                  # inches float for calculations
CONTENT_MAX_HEIGHT = SLIDE_HEIGHT - SAFE_TOP - SAFE_BOTTOM    # ~6.55" (EMU)

# ── Vertical rhythm (slide = 7.5" tall) ──
#   y=0.30"  Section title (small gray text)
#   y=0.75"  Accent gold line (3pt)
#   y=1.10"  Slide title start (dynamic height)
#   y=?      Body content start (determined by draw_slide_title return)
#   y=6.15"  Content safe bottom (0.4" buffer before footer)
#   y=6.60"  Footer divider line (1pt)
#   y=6.75"  Page number
#   y=7.50"  Slide edge

HEADER_Y = Inches(0.3)
HEADER_LINE_Y = Inches(0.75)
TITLE_Y = Inches(1.1)
BODY_Y = Inches(1.85)
CONTENT_BOTTOM = Inches(6.15)
FOOTER_LINE_Y = Inches(6.6)
FOOTER_TEXT_Y = Inches(6.75)

# ═══════════════════════════════════════════════════════════════════
# Common drawing helpers
# ═══════════════════════════════════════════════════════════════════


def _set_font(run, name: str = BODY_FONT, size=Pt(18), color=TEXT_MAIN, bold: bool = False):
    """Apply font settings to a run."""
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(
    slide: PptxSlide,
    left,
    top,
    width,
    height,
    text: str = "",
    font_name: str = BODY_FONT,
    font_size=Pt(18),
    color=TEXT_MAIN,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
    word_wrap: bool = True,
):
    """Add a single-paragraph textbox and return (shape, text_frame, paragraph)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    _set_font(p.runs[0] if p.runs else p.add_run(), font_name, font_size, color, bold)
    return tb, tf, p


def add_paragraph(tf, text: str, font_name=BODY_FONT, font_size=Pt(18), color=TEXT_MAIN, bold=False, space_after=Pt(8), alignment=PP_ALIGN.LEFT):
    """Append a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.alignment = alignment
    p.space_after = space_after
    _set_font(p.runs[0] if p.runs else p.add_run(), font_name, font_size, color, bold)
    return p


def add_bullet_paragraph(tf, text: str, font_size=Pt(18), color=TEXT_MAIN, indent_level: int = 0, space_after=Pt(6)):
    """Add a bullet-point paragraph with proper indentation."""
    p = tf.add_paragraph()
    p.text = text
    p.level = indent_level
    p.space_after = space_after
    _set_font(p.runs[0] if p.runs else p.add_run(), BODY_FONT, font_size, color)
    return p


def draw_header(slide: PptxSlide, section_title: Optional[str] = None):
    """Draw the standard slide header: thin gold accent line + optional section title.

    Returns the Y position below the header (where the slide title should go).
    """
    # Thin gold accent line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        SAFE_LEFT, HEADER_LINE_Y,
        Inches(2.5), Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if section_title:
        add_textbox(
            slide,
            SAFE_LEFT, HEADER_Y, CONTENT_WIDTH, Inches(0.4),
            text=section_title,
            font_size=SIZE_CAPTION,
            color=TEXT_MUTED,
        )

    return TITLE_Y  # slide title should start at TITLE_Y


def draw_footer(slide: PptxSlide, page_number: int, total_pages: int):
    """Draw the standard slide footer with page number."""
    # Thin separator line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        SAFE_LEFT, FOOTER_LINE_Y,
        CONTENT_WIDTH, Pt(1),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER
    line.line.fill.background()

    # Page number (left-aligned)
    add_textbox(
        slide,
        SAFE_LEFT, FOOTER_TEXT_Y,
        Inches(3), Inches(0.3),
        text=f"{page_number} / {total_pages}",
        font_size=SIZE_CAPTION,
        color=TEXT_MUTED,
    )


def draw_rounded_rect(slide: PptxSlide, left, top, width, height, fill_color=LIGHT_BG, border_color=None):
    """Draw a rounded rectangle shape (used for cards, panels)."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def draw_decorative_circle(slide: PptxSlide, left, top, diameter, color=ACCENT):
    """Draw a circle for decorative purposes."""
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        left, top, diameter, diameter,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def truncate_text(text: str, max_chars: int) -> str:
    """Truncate text to a maximum character count, adding ellipsis if needed."""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1] + "…"


def calculate_text_rows(text: str, font_size_pt: float, max_width_inches: float) -> int:
    """Roughly estimate how many rows text will occupy.

    CJK characters are full-width (~1.0× font_size). Mixed CJK/Western text
    averages ~0.85× font_size per character. We use a conservative 0.90
    to ensure we never underestimate rows.
    """
    char_width_pt = font_size_pt * 0.90  # conservative CJK-majority estimate
    chars_per_line = max(1, int(max_width_inches * 72 / char_width_pt))
    lines = (len(text) + chars_per_line - 1) // chars_per_line
    return max(1, lines)


def estimate_text_height(text: str, font_size_pt: float, max_width_inches: float) -> int:
    """Estimate the EMU height needed for a text box.

    Uses 1.65× line-height and a 25% safety factor to ensure CJK text
    never overflows vertically in PowerPoint.
    """
    rows = calculate_text_rows(text, font_size_pt, max_width_inches)
    line_height_inches = font_size_pt / 72 * 1.65
    safe_height = rows * line_height_inches * 1.25  # 25% extra for safety
    return Inches(safe_height)


def draw_slide_title(
    slide: PptxSlide,
    title: str,
    y: int = TITLE_Y,
    font_size=SIZE_SLIDE_TITLE,
    color=PRIMARY,
) -> int:
    """Draw the slide title at the given Y position.

    Returns the Y position immediately below the title (for body content).
    The textbox height is calculated dynamically based on text length and
    font size to prevent overflow when the title wraps to multiple lines.
    """
    text_height = estimate_text_height(title, font_size.pt, CONTENT_WIDTH_INCHES)
    # Minimum 1.0" for single line, more for multi-line titles
    box_height = max(text_height, Inches(1.0))

    add_textbox(
        slide,
        SAFE_LEFT, y, CONTENT_WIDTH, box_height,
        text=title,
        font_size=font_size,
        color=color,
        bold=True,
    )
    return y + box_height + Inches(0.4)  # ample gap before body content


def draw_bullet_body(
    slide: PptxSlide,
    bullets: list,
    start_y: int,
    max_y: int = CONTENT_BOTTOM,
    font_size=SIZE_BODY,
    color=TEXT_MAIN,
    max_bullets: int = 5,
) -> int:
    """Draw bullet points between start_y and max_y.

    Returns the Y position after the last bullet (where next content
    element should go). If bullets would exceed max_y, they get a
    shorter textbox to prevent overlap.
    """
    available = max_y - start_y
    if available < Inches(0.5):
        return start_y  # not enough room — skip

    body_box = slide.shapes.add_textbox(
        SAFE_LEFT, start_y, CONTENT_WIDTH, available,
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:max_bullets]):
        bullet = truncate_text(bullet, 200)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = BODY_FONT
        p.space_after = Pt(10)

    # Estimate where bullets end: each bullet ≈ 1.2 lines
    count = min(len(bullets), max_bullets)
    line_h = font_size.pt / 72 * 1.65
    estimated_height = Inches(count * 1.2 * line_h * 1.25)
    end_y = start_y + estimated_height + Inches(0.1)

    # Never exceed max_y
    return min(end_y, max_y - Inches(0.05))


def fit_font_to_width(
    text: str,
    max_width_inches: float,
    start_size_pt: float,
    min_size_pt: float = 10.0,
    char_factor: float = 0.92,
) -> float:
    """Find the largest font size (≤ start_size_pt) that keeps text on one line.

    If text fits at start_size_pt, return that. Otherwise shrink in 1pt
    increments until the text fits within max_width_inches. Never drops
    below min_size_pt.
    """
    size = start_size_pt
    while size > min_size_pt:
        char_width_pt = size * char_factor
        chars_per_line = max(1, int(max_width_inches * 72 / char_width_pt))
        if len(text) <= chars_per_line:
            return size
        size -= 1.0
    return min_size_pt


def add_fitted_textbox(
    slide: PptxSlide,
    left,
    top,
    width,
    height,
    text: str,
    start_font_size,
    min_font_size=Pt(10),
    font_name: str = BODY_FONT,
    color=TEXT_MAIN,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
) -> float:
    """Add a textbox with auto-shrunk font size to keep text on one line.

    Returns the actual font size used (in Pt). If the text needs to shrink
    below min_font_size, it is truncated with an ellipsis instead.
    """
    width_inches = width / 914400.0  # EMU → inches conversion

    actual_pt = fit_font_to_width(
        text, width_inches, start_font_size.pt, min_font_size.pt
    )
    actual_size = Pt(int(actual_pt))

    # If at minimum size and still doesn't fit, truncate
    display_text = text
    if actual_pt <= min_font_size.pt:
        display_text = truncate_text(text, int(width_inches * 72 / (min_font_size.pt * 0.92)))

    add_textbox(
        slide, left, top, width, height,
        text=display_text,
        font_size=actual_size,
        color=color,
        bold=bold,
        font_name=font_name,
        alignment=alignment,
    )
    return actual_size
