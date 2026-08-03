"""Cover / title slide layout with decorative background elements."""

from pptx.dml.color import RGBColor
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from PPT_Generator.design import (
    ACCENT,
    BODY_Y,
    CONTENT_WIDTH,
    HEADER_LINE_Y,
    PRIMARY,
    SAFE_LEFT,
    SECONDARY,
    SIZE_SUBTITLE,
    SIZE_TITLE_COVER,
    TEXT_MUTED,
    WHITE,
    add_paragraph,
    add_textbox,
    draw_decorative_circle,
)
from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout


class TitleSlideLayout(BaseLayout):
    layout_id = "title_slide"

    def render(self, slide: Slide, prs_slide: PptxSlide, total_pages: int = 0) -> None:
        # ── Background: navy block covering left 60% ──
        bg_shape = prs_slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(0), Inches(0),
            Inches(8), Inches(7.5),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = PRIMARY
        bg_shape.line.fill.background()

        # ── Decorative gold accent bar ──
        bar = prs_slide.shapes.add_shape(
            1,
            Inches(8), Inches(0),
            Pt(6), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # ── Decorative circles in the background ──
        draw_decorative_circle(prs_slide, Inches(5.5), Inches(4.5), Inches(2.5), SECONDARY)

        # ── Title on the left (dark area) ──
        add_textbox(
            prs_slide, Inches(1.0), Inches(2.2), Inches(6.5), Inches(2.0),
            text=slide.title,
            font_size=SIZE_TITLE_COVER,
            color=WHITE,
            bold=True,
        )

        # ── Decorative line under title ──
        line = prs_slide.shapes.add_shape(
            1, Inches(1.0), Inches(4.5), Inches(2.0), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # ── Subtitle ──
        if slide.subtitle:
            add_textbox(
                prs_slide, Inches(1.0), Inches(4.8), Inches(6.0), Inches(1.2),
                text=slide.subtitle,
                font_size=SIZE_SUBTITLE,
                color=RGBColor(0xCC, 0xCC, 0xCC),
            )

        # ── Right-side metadata area ──
        if slide.source_notes:
            # Show source/topic info on the right white area
            tb, _, _ = add_textbox(
                prs_slide, Inches(8.8), Inches(5.5), Inches(3.8), Inches(1.5),
                text="",
                font_size=Pt(11),
                color=TEXT_MUTED,
            )
            for i, note in enumerate(slide.source_notes[:3]):
                if i == 0:
                    tb.text_frame.paragraphs[0].text = note
                else:
                    add_paragraph(tb.text_frame, note, font_size=Pt(11), color=TEXT_MUTED)
