from pptx.util import Inches, Pt

from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout
from PPT_Generator.templates.styles import COLORS, FONTS


class TitleLayout(BaseLayout):
    layout_id = "title"

    def render(self, slide: Slide, prs_slide) -> None:
        title_box = prs_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["chinese"]

        if slide.subtitle:
            sub_box = prs_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide.subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS["muted"]
            p.font.name = FONTS["chinese"]
