"""Two-column comparison slide layout."""

from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from PPT_Generator.design import (
    BODY_FONT,
    CONTENT_BOTTOM,
    CONTENT_WIDTH,
    LIGHT_BG,
    SAFE_LEFT,
    TEXT_MAIN,
    draw_footer,
    draw_header,
    draw_slide_title,
    truncate_text,
)
from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout

COL_WIDTH = Inches(5.3)
COL_GAP = Inches(0.3)


class TwoColumnLayout(BaseLayout):
    layout_id = "two_column"

    def render(self, slide: Slide, prs_slide: PptxSlide, total_pages: int = 0) -> None:
        draw_header(prs_slide, section_title=slide.section_title)
        draw_footer(prs_slide, slide.page_number, total_pages)

        # ── Slide title (dynamic height) ──
        col_top = draw_slide_title(prs_slide, slide.title)
        col_height = CONTENT_BOTTOM - col_top

        # ── Left column ──
        left_x = SAFE_LEFT
        left_bg = prs_slide.shapes.add_shape(
            5, left_x, col_top, COL_WIDTH, col_height,
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = LIGHT_BG
        left_bg.line.fill.background()

        left_tf = prs_slide.shapes.add_textbox(
            left_x + Inches(0.3), col_top + Inches(0.2),
            COL_WIDTH - Inches(0.6), col_height - Inches(0.4),
        ).text_frame
        left_tf.word_wrap = True
        for i, item in enumerate(slide.left_column[:5]):
            item = truncate_text(item, 150)
            p = left_tf.paragraphs[0] if i == 0 else left_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = BODY_FONT
            p.space_after = Pt(8)

        # ── Right column ──
        right_x = SAFE_LEFT + COL_WIDTH + COL_GAP
        right_bg = prs_slide.shapes.add_shape(
            5, right_x, col_top, COL_WIDTH, col_height,
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = LIGHT_BG
        right_bg.line.fill.background()

        right_tf = prs_slide.shapes.add_textbox(
            right_x + Inches(0.3), col_top + Inches(0.2),
            COL_WIDTH - Inches(0.6), col_height - Inches(0.4),
        ).text_frame
        right_tf.word_wrap = True
        for i, item in enumerate(slide.right_column[:5]):
            item = truncate_text(item, 150)
            p = right_tf.paragraphs[0] if i == 0 else right_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = BODY_FONT
            p.space_after = Pt(8)
