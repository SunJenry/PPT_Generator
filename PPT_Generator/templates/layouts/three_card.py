"""Three-card layout — three parallel cards side by side."""

from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from PPT_Generator.design import (
    ACCENT,
    BODY_FONT,
    CONTENT_BOTTOM,
    CONTENT_WIDTH,
    LIGHT_BG,
    PRIMARY,
    SAFE_LEFT,
    TEXT_MAIN,
    add_textbox,
    draw_footer,
    draw_header,
    draw_slide_title,
    truncate_text,
)
from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout

CARD_WIDTH = Inches(3.5)
CARD_GAP = Inches(0.35)


class ThreeCardLayout(BaseLayout):
    layout_id = "three_card"

    def render(self, slide: Slide, prs_slide: PptxSlide, total_pages: int = 0) -> None:
        draw_header(prs_slide, section_title=slide.section_title)
        draw_footer(prs_slide, slide.page_number, total_pages)

        # ── Slide title (dynamic height) ──
        card_top = draw_slide_title(prs_slide, slide.title)
        card_height = CONTENT_BOTTOM - card_top

        # ── Draw 3 cards ──
        cards = slide.cards[:3]
        for idx, card in enumerate(cards):
            card_left = SAFE_LEFT + idx * (CARD_WIDTH + CARD_GAP)

            # Card background
            card_shape = prs_slide.shapes.add_shape(
                5, card_left, card_top, CARD_WIDTH, card_height,
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = LIGHT_BG
            card_shape.line.fill.background()

            # Accent line on top of card
            accent_bar = prs_slide.shapes.add_shape(
                1, card_left, card_top, CARD_WIDTH, Pt(4),
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = ACCENT
            accent_bar.line.fill.background()

            # Card title
            card_title = card.get("title", f"要点 {idx + 1}")
            add_textbox(
                prs_slide,
                card_left + Inches(0.25), card_top + Inches(0.3),
                CARD_WIDTH - Inches(0.5), Inches(0.6),
                text=truncate_text(card_title, 30),
                font_size=Pt(20), color=PRIMARY, bold=True,
            )

            # Card body
            card_body = card.get("body", "")
            body_tb = prs_slide.shapes.add_textbox(
                card_left + Inches(0.25), card_top + Inches(1.1),
                CARD_WIDTH - Inches(0.5), card_height - Inches(1.5),
            )
            body_tb.text_frame.word_wrap = True
            p = body_tb.text_frame.paragraphs[0]
            p.text = truncate_text(card_body, 200)
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = BODY_FONT
            p.space_after = Pt(6)
