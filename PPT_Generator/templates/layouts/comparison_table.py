"""Comparison table layout — multi-item side-by-side comparison."""

from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from PPT_Generator.design import (
    BODY_FONT,
    CONTENT_BOTTOM,
    CONTENT_WIDTH,
    LIGHT_BG,
    PRIMARY,
    SAFE_LEFT,
    TEXT_MAIN,
    WHITE,
    add_textbox,
    draw_footer,
    draw_header,
    draw_slide_title,
    truncate_text,
)
from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout


class ComparisonTableLayout(BaseLayout):
    layout_id = "comparison_table"

    def render(self, slide: Slide, prs_slide: PptxSlide, total_pages: int = 0) -> None:
        draw_header(prs_slide, section_title=slide.section_title)
        draw_footer(prs_slide, slide.page_number, total_pages)

        # ── Slide title (dynamic height) ──
        table_top = draw_slide_title(prs_slide, slide.title)

        table_data = slide.table
        if not table_data or len(table_data) < 2:
            return

        rows = len(table_data)
        cols = len(table_data[0])
        if cols == 0:
            return

        col_width = CONTENT_WIDTH / cols
        row_height = min(Inches(0.65), (CONTENT_BOTTOM - table_top) / rows)

        for r, row in enumerate(table_data[:8]):  # max 8 rows
            for c, cell_text in enumerate(row[:6]):  # max 6 cols
                cell_left = SAFE_LEFT + col_width * c
                cell_top = table_top + row_height * r

                # Cell background
                if r == 0:
                    fill_color = PRIMARY
                    text_color = WHITE
                elif r % 2 == 0:
                    fill_color = LIGHT_BG
                    text_color = TEXT_MAIN
                else:
                    fill_color = WHITE
                    text_color = TEXT_MAIN

                cell_shape = prs_slide.shapes.add_shape(
                    1, cell_left, cell_top, col_width, row_height,
                )
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = fill_color
                cell_shape.line.color.rgb = LIGHT_BG
                cell_shape.line.width = Pt(0.5)

                # Cell text
                cell_tb = prs_slide.shapes.add_textbox(
                    cell_left + Inches(0.1), cell_top + Inches(0.05),
                    col_width - Inches(0.2), row_height - Inches(0.1),
                )
                cell_tb.text_frame.word_wrap = True
                p = cell_tb.text_frame.paragraphs[0]
                p.text = truncate_text(cell_text, 40)
                p.font.size = Pt(12) if r > 0 else Pt(13)
                p.font.color.rgb = text_color
                p.font.name = BODY_FONT
                p.font.bold = (r == 0)
