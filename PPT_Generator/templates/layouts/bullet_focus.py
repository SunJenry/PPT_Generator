from pptx.util import Inches, Pt

from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout
from PPT_Generator.templates.styles import COLORS, FONTS


class BulletFocusLayout(BaseLayout):
    layout_id = "bullet_focus"

    def render(self, slide: Slide, prs_slide) -> None:
        title_box = prs_slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["chinese"]

        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.733)
        height = Inches(5)
        body_box = prs_slide.shapes.add_textbox(left, top, width, height)
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(slide.bullets[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS["text"]
            p.font.name = FONTS["chinese"]
            p.space_after = Pt(12)
