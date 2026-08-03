"""Timeline layout — horizontal chronological sequence."""

from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from PPT_Generator.design import (
    ACCENT,
    BODY_FONT,
    CONTENT_BOTTOM,
    CONTENT_WIDTH,
    PRIMARY,
    SAFE_LEFT,
    SECONDARY,
    TEXT_MAIN,
    add_textbox,
    draw_footer,
    draw_header,
    draw_slide_title,
    truncate_text,
)
from PPT_Generator.models import Slide
from PPT_Generator.templates.base import BaseLayout

MAX_ITEMS = 6
DOT_SIZE = Inches(0.2)


class TimelineLayout(BaseLayout):
    layout_id = "timeline"

    def render(self, slide: Slide, prs_slide: PptxSlide, total_pages: int = 0) -> None:
        draw_header(prs_slide, section_title=slide.section_title)
        draw_footer(prs_slide, slide.page_number, total_pages)

        # ── Slide title (dynamic height) ──
        body_y = draw_slide_title(prs_slide, slide.title)

        items = slide.timeline_items[:MAX_ITEMS]
        if not items:
            return

        n = len(items)
        usable_width = CONTENT_WIDTH - Inches(1.0)
        spacing = usable_width / max(n - 1, 1)

        # Timeline line centered in the remaining content area
        available_height = CONTENT_BOTTOM - body_y
        line_y = body_y + available_height / 2

        # ── Horizontal line ──
        line = prs_slide.shapes.add_shape(
            1,
            SAFE_LEFT + Inches(0.5), line_y,
            usable_width, Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = SECONDARY
        line.line.fill.background()

        # ── Dots and labels ──
        for idx, item in enumerate(items):
            x_center = SAFE_LEFT + Inches(0.5) + spacing * idx

            # Dot
            dot = prs_slide.shapes.add_shape(
                9,  # OVAL
                x_center - DOT_SIZE // 2, line_y - DOT_SIZE // 2,
                DOT_SIZE, DOT_SIZE,
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT
            dot.line.fill.background()

            # Date (above the line)
            date_text = item.get("date", "")
            add_textbox(
                prs_slide,
                x_center - Inches(0.8), line_y - Inches(0.55),
                Inches(1.6), Inches(0.4),
                text=truncate_text(date_text, 20),
                font_size=Pt(12), color=ACCENT, bold=True,
            )

            # Event (below the line)
            event_text = item.get("event", "")
            event_box = prs_slide.shapes.add_textbox(
                x_center - Inches(0.8), line_y + Inches(0.15),
                Inches(1.6), Inches(1.5),
            )
            event_box.text_frame.word_wrap = True
            p = event_box.text_frame.paragraphs[0]
            p.text = truncate_text(event_text, 60)
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = BODY_FONT
