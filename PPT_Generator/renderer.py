import os
import sys
import tempfile
from typing import Optional

import httpx
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from PPT_Generator.design import SLIDE_HEIGHT, SLIDE_WIDTH, check_and_report_overlaps
from PPT_Generator.image_search import ImageSearchClient
from PPT_Generator.models import Presentation, Slide
from PPT_Generator.templates.registry import TemplateRegistry


class Renderer:
    def __init__(self, templates: TemplateRegistry, image_client: Optional[ImageSearchClient] = None):
        self.templates = templates
        self.image_client = image_client

    def render(self, presentation: Presentation, output_path: str) -> None:
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]

        total = presentation.total_pages
        overlap_count = 0
        for i, slide_model in enumerate(presentation.slides, start=1):
            prs_slide = prs.slides.add_slide(blank_layout)

            # Ensure page_number is correct and pass total for footer
            slide_model.page_number = i

            layout = self.templates.get(slide_model.layout_id)
            layout.render(slide_model, prs_slide, total_pages=total)

            if slide_model.image_keyword and self.image_client:
                self._add_image(prs_slide, slide_model)

            # Check for content overlap
            n = check_and_report_overlaps(prs_slide, i)
            overlap_count += n

            # Print progress every 5 slides
            if i % 5 == 0 or i == total:
                print(f"       rendering slide {i}/{total}...")

        if overlap_count > 0:
            print(f"       ⚠  {overlap_count} overlap(s) detected across all slides", file=sys.stderr)

        prs.save(output_path)

    def _add_image(self, prs_slide, slide_model: Slide) -> None:
        tmp_path = None
        try:
            url = self.image_client.search(slide_model.image_keyword)
            if url:
                with httpx.Client(timeout=20.0) as client:
                    response = client.get(url)
                    response.raise_for_status()
                    with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as f:
                        f.write(response.content)
                        tmp_path = f.name
                prs_slide.shapes.add_picture(tmp_path, Inches(9), Inches(0.6), width=Inches(3.5))
        except Exception:
            print(f"Warning: failed to add image for keyword '{slide_model.image_keyword}'", file=sys.stderr)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)
